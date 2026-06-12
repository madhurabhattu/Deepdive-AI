"""
DeepDive AI — PowerPoint Presentation Generator

Generates a premium dark-themed 7-slide PPTX from a ResearchReport using python-pptx.
"""

from __future__ import annotations

import logging
import re
from datetime import datetime
from pathlib import Path
from typing import cast

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from utils.localization import get_text
from utils.report_schema import ResearchReport

logger = logging.getLogger(__name__)

# ── Constants ───────────────────────────────────────────────────────
OUTPUT_DIR = Path(__file__).resolve().parent.parent / "output"

# Premium SaaS Dark Palette
BG_COLOR = RGBColor(0x0B, 0x0F, 0x19)  # #0B0F19 - Main background
CARD_BG_COLOR = RGBColor(0x1A, 0x22, 0x36)  # #1A2236 - Card background
ACCENT_VIOLET = RGBColor(0x7C, 0x3A, 0xED)  # #7C3AED - Accent 1 (Violet)
ACCENT_PURPLE = RGBColor(0xA8, 0x55, 0xF7)  # #A855F7 - Accent 2 (Purple)
ACCENT_PINK = RGBColor(0xEC, 0x48, 0x99)  # #EC4899 - Accent 3 (Pink)
TEXT_PRIMARY = RGBColor(0xF8, 0xFA, 0xFC)  # #F8FAFC - Main text
TEXT_SECONDARY = RGBColor(0x94, 0xA3, 0xB8)  # #94A3B8 - Secondary text
BORDER_COLOR = RGBColor(0x2E, 0x3A, 0x56)  # Card border outline


def sanitise_filename(topic: str) -> str:
    """Replace non-alphanumeric characters with underscores and truncate."""
    clean = re.sub(r"[^a-zA-Z0-9]+", "_", topic).strip("_").lower()
    return clean[:80] if len(clean) > 80 else clean


def _set_slide_bg(slide, color: RGBColor) -> None:
    """Set a solid background colour for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _get_font_name(lang: str = "en") -> str:
    """Return appropriate font name based on language selection."""
    if lang in ["hi", "mr", "te"]:
        return "Nirmala UI"
    return "Outfit"


def _add_card(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    bg_color: RGBColor = CARD_BG_COLOR,
    border_color: RGBColor = BORDER_COLOR,
):
    """Add a card container with a solid background and custom border."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def _add_textbox(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    text: str,
    font_size: int = 14,
    bold: bool = False,
    color: RGBColor = TEXT_PRIMARY,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
    line_spacing: int = 0,
    lang: str = "en",
) -> None:
    """Add a simple text box to a slide with language-compatible typography."""
    tx_box = slide.shapes.add_textbox(left, top, width, height)
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = _get_font_name(lang)
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    if line_spacing > 0:
        p.line_spacing = Pt(line_spacing)
    return tx_box


def _create_slide_with_header(prs, number_str: str, title: str, lang: str = "en"):
    """Create a new slide with dark background and standard header layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    _set_slide_bg(slide, BG_COLOR)

    # Section indicator
    _add_textbox(
        slide,
        Inches(0.8),
        Inches(0.4),
        Inches(11.7),
        Inches(0.4),
        number_str.upper(),
        font_size=11,
        bold=True,
        color=ACCENT_PINK,
        lang=lang,
    )
    # Slide Title
    _add_textbox(
        slide,
        Inches(0.8),
        Inches(0.7),
        Inches(11.7),
        Inches(0.6),
        title,
        font_size=24,
        bold=True,
        color=TEXT_PRIMARY,
        lang=lang,
    )

    # Thin divider line below title
    connector = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(11.733), Inches(0.02)
    )
    connector.fill.solid()
    connector.fill.fore_color.rgb = BORDER_COLOR
    connector.line.fill.background()

    return slide


def build_ppt(report: ResearchReport, lang: str = "en") -> str:
    """Generate a premium dark-themed 7-slide presentation in target language."""
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    filename = f"{sanitise_filename(report.topic)}_presentation.pptx"
    filepath = OUTPUT_DIR / filename
    date_str = datetime.now().strftime("%B %d, %Y")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Slide 0: Cover Slide (Title Slide) ───────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG_COLOR)

    # Asymmetric vertical color block on the left
    vertical_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(0.1), Inches(4.0)
    )
    vertical_bar.fill.solid()
    vertical_bar.fill.fore_color.rgb = ACCENT_VIOLET
    vertical_bar.line.fill.background()

    # Title content
    _add_textbox(
        slide,
        Inches(1.2),
        Inches(1.8),
        Inches(11),
        Inches(0.5),
        get_text("app_title", lang).upper(),
        font_size=16,
        bold=True,
        color=ACCENT_PINK,
        lang=lang,
    )
    _add_textbox(
        slide,
        Inches(1.2),
        Inches(2.4),
        Inches(11),
        Inches(2.2),
        report.topic,
        font_size=38,
        bold=True,
        color=TEXT_PRIMARY,
        lang=lang,
    )
    _add_textbox(
        slide,
        Inches(1.2),
        Inches(4.7),
        Inches(11),
        Inches(0.4),
        get_text("app_tagline", lang),
        font_size=16,
        color=TEXT_SECONDARY,
        lang=lang,
    )
    _add_textbox(
        slide,
        Inches(1.2),
        Inches(5.2),
        Inches(11),
        Inches(0.4),
        f"{get_text('built_with', lang)} DeepDive AI — {date_str}",
        font_size=12,
        color=TEXT_SECONDARY,
        lang=lang,
    )

    # ── Slide 1: Executive Summary ───────────────────────────────────
    slide = _create_slide_with_header(
        prs,
        get_text("sec_summary", lang),
        get_text("strategic_overview_heading", lang),
        lang=lang,
    )

    # Left Card: Summary Text
    _add_card(slide, Inches(0.8), Inches(1.8), Inches(7.5), Inches(4.8))
    _add_textbox(
        slide,
        Inches(1.1),
        Inches(2.1),
        Inches(6.9),
        Inches(4.2),
        report.executive_summary,
        font_size=15,
        color=TEXT_PRIMARY,
        line_spacing=22,
        lang=lang,
    )

    # Right Card: Key Metrics Highlight
    _add_card(slide, Inches(8.6), Inches(1.8), Inches(3.933), Inches(4.8))
    _add_textbox(
        slide,
        Inches(8.9),
        Inches(2.1),
        Inches(3.3),
        Inches(0.4),
        get_text("primary_metrics_label", lang),
        font_size=12,
        bold=True,
        color=ACCENT_PINK,
        lang=lang,
    )

    # Place up to 3 metrics in the right column
    metrics = report.statistics[:3]
    for idx, metric in enumerate(metrics):
        top_offset: int = Inches(2.6 + idx * 1.3)
        _add_textbox(
            slide,
            Inches(8.9),
            top_offset,
            Inches(3.3),
            Inches(0.5),
            metric.get("value", "N/A"),
            font_size=28,
            bold=True,
            color=ACCENT_VIOLET,
            lang=lang,
        )
        _add_textbox(
            slide,
            Inches(8.9),
            top_offset + Inches(0.5),
            Inches(3.3),
            Inches(0.4),
            metric.get("label", "Metric").upper(),
            font_size=10,
            bold=True,
            color=TEXT_SECONDARY,
            lang=lang,
        )

    # ── Slide 2: Background & Context ───────────────────────────────
    slide = _create_slide_with_header(
        prs,
        get_text("sec_background", lang),
        get_text("why_this_matters_heading", lang),
        lang=lang,
    )

    # Left Card: Context Text
    _add_card(slide, Inches(0.8), Inches(1.8), Inches(7.5), Inches(4.8))
    _add_textbox(
        slide,
        Inches(1.1),
        Inches(2.1),
        Inches(6.9),
        Inches(4.2),
        report.background_context,
        font_size=15,
        color=TEXT_PRIMARY,
        line_spacing=22,
        lang=lang,
    )

    # Right Card: Relevance Highlight
    _add_card(slide, Inches(8.6), Inches(1.8), Inches(3.933), Inches(4.8))
    _add_textbox(
        slide,
        Inches(8.9),
        Inches(2.1),
        Inches(3.3),
        Inches(0.4),
        get_text("landscape_importance_label", lang),
        font_size=12,
        bold=True,
        color=ACCENT_PINK,
        lang=lang,
    )

    # Context checklist or narrative summary
    _add_textbox(
        slide,
        Inches(8.9),
        Inches(2.6),
        Inches(3.3),
        Inches(3.8),
        get_text("why_it_matters_text", lang),
        font_size=15,
        color=TEXT_PRIMARY,
        line_spacing=22,
        lang=lang,
    )

    # ── Slide 3: Core Concepts ───────────────────────────────────────
    slide = _create_slide_with_header(
        prs,
        get_text("sec_concepts", lang),
        get_text("foundational_principles_heading", lang),
        lang=lang,
    )

    # 3-Column Concept Cards
    col_width = Inches(3.64)
    for idx, concept in enumerate(report.core_concepts[:3]):
        left_pos: int = Inches(0.8 + idx * 4.04)
        _add_card(slide, left_pos, Inches(1.8), col_width, Inches(4.8))

        # Concept Index
        _add_textbox(
            slide,
            left_pos + Inches(0.3),
            Inches(2.1),
            col_width - Inches(0.6),
            Inches(0.3),
            f"{get_text('concept_label', lang)} 0{idx + 1}",
            font_size=11,
            bold=True,
            color=ACCENT_PINK,
            lang=lang,
        )
        # Term
        _add_textbox(
            slide,
            left_pos + Inches(0.3),
            Inches(2.5),
            col_width - Inches(0.6),
            Inches(0.8),
            concept.get("term", "Term"),
            font_size=18,
            bold=True,
            color=TEXT_PRIMARY,
            lang=lang,
        )
        # Definition
        _add_textbox(
            slide,
            left_pos + Inches(0.3),
            Inches(3.5),
            col_width - Inches(0.6),
            Inches(2.8),
            concept.get("definition", ""),
            font_size=13,
            color=TEXT_SECONDARY,
            line_spacing=18,
            lang=lang,
        )

    # ── Slide 4: Research Findings & Analysis ─────────────────────────
    slide = _create_slide_with_header(
        prs,
        get_text("sec_insights", lang),
        get_text("findings_analysis_heading", lang),
        lang=lang,
    )

    # Left Card: Vertical stack of Key Insights
    _add_card(slide, Inches(0.8), Inches(1.8), Inches(7.5), Inches(4.8))
    _add_textbox(
        slide,
        Inches(1.1),
        Inches(2.1),
        Inches(6.9),
        Inches(0.4),
        get_text("key_analysis_insights", lang),
        font_size=12,
        bold=True,
        color=ACCENT_PINK,
        lang=lang,
    )

    tx_box = slide.shapes.add_textbox(
        Inches(1.1), Inches(2.6), Inches(6.9), Inches(3.8)
    )
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    for idx, insight in enumerate(report.key_insights[:4]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"▸  {insight}"
        p.font.name = _get_font_name(lang)
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(12)
        p.line_spacing = Pt(18)

    # Right Card: Statistics Stack
    _add_card(slide, Inches(8.6), Inches(1.8), Inches(3.933), Inches(4.8))
    _add_textbox(
        slide,
        Inches(8.9),
        Inches(2.1),
        Inches(3.3),
        Inches(0.4),
        get_text("supporting_data", lang),
        font_size=12,
        bold=True,
        color=ACCENT_PINK,
        lang=lang,
    )

    stats_list = (
        report.statistics[3:6]
        if len(report.statistics) >= 6
        else report.statistics[:3]
    )
    for idx, stat in enumerate(stats_list):
        top_offset = Inches(2.6 + idx * 1.3)
        _add_textbox(
            slide,
            Inches(8.9),
            top_offset,
            Inches(3.3),
            Inches(0.5),
            stat.get("value", "N/A"),
            font_size=26,
            bold=True,
            color=ACCENT_PURPLE,
            lang=lang,
        )
        _add_textbox(
            slide,
            Inches(8.9),
            top_offset + Inches(0.45),
            Inches(3.3),
            Inches(0.4),
            stat.get("label", "Metric").upper(),
            font_size=10,
            bold=True,
            color=TEXT_SECONDARY,
            lang=lang,
        )

    # ── Slide 5: Benefits, Challenges & Risks ─────────────────────────
    slide = _create_slide_with_header(
        prs,
        get_text("sec_benefits", lang),
        get_text("advantages_constraints_heading", lang),
        lang=lang,
    )

    # 3 columns mapping Benefits, Challenges, Risks
    col_width = Inches(3.64)

    # Color mappings for the columns
    green_accent = RGBColor(0x4A, 0xDE, 0x80)
    orange_accent = RGBColor(0xFB, 0xBF, 0x24)
    rose_accent = RGBColor(0xF8, 0x71, 0x71)

    columns_data = [
        {
            "title": get_text("benefit_col_title", lang),
            "type": "benefit",
            "color": green_accent,
            "left": Inches(0.8),
        },
        {
            "title": get_text("challenge_col_title", lang),
            "type": "challenge",
            "color": orange_accent,
            "left": Inches(4.84),
        },
        {
            "title": get_text("risk_col_title", lang),
            "type": "risk",
            "color": rose_accent,
            "left": Inches(8.88),
        },
    ]

    for col in columns_data:
        col_left: int = cast(int, col["left"])
        col_color: RGBColor = cast(RGBColor, col["color"])
        col_title: str = cast(str, col["title"])
        col_type: str = cast(str, col["type"])

        _add_card(slide, col_left, Inches(1.8), col_width, Inches(4.8))
        _add_textbox(
            slide,
            col_left + Inches(0.3),
            Inches(2.1),
            col_width - Inches(0.6),
            Inches(0.4),
            col_title,
            font_size=11,
            bold=True,
            color=col_color,
            lang=lang,
        )

        # Populate matching items
        items = [
            i
            for i in report.benefits_challenges_risks
            if i.get("type", "").lower() == col_type
        ]
        if not items:
            items = [
                i
                for i in report.benefits_challenges_risks
                if col_type in i.get("type", "").lower()
            ]

        tx_box = slide.shapes.add_textbox(
            col_left + Inches(0.3),
            Inches(2.6),
            col_width - Inches(0.6),
            Inches(3.8),
        )
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)

        for idx, it in enumerate(items[:3]):
            if idx == 0 and len(tf.paragraphs[0].text) == 0:
                p_title = tf.paragraphs[0]
            else:
                p_title = tf.add_paragraph()
            p_title.text = f"▸  {it.get('item', 'Item')}"
            p_title.font.name = _get_font_name(lang)
            p_title.font.size = Pt(13)
            p_title.font.bold = True
            p_title.font.color.rgb = TEXT_PRIMARY
            p_title.space_after = Pt(2)
            p_title.line_spacing = Pt(16)

            if it.get("description"):
                p_desc = tf.add_paragraph()
                p_desc.text = it.get("description", "")
                p_desc.font.name = _get_font_name(lang)
                p_desc.font.size = Pt(11)
                p_desc.font.color.rgb = TEXT_SECONDARY
                p_desc.left_indent = Inches(0.25)
                p_desc.space_after = Pt(12)
                p_desc.line_spacing = Pt(15)

    # ── Slide 6: Real-World Applications ─────────────────────────────
    slide = _create_slide_with_header(
        prs,
        get_text("sec_apps", lang),
        get_text("practical_adoption_heading", lang),
        lang=lang,
    )

    col_width = Inches(3.64)
    for idx, app in enumerate(report.real_world_applications[:3]):
        left_pos = Inches(0.8 + idx * 4.04)
        _add_card(slide, left_pos, Inches(1.8), col_width, Inches(4.8))

        # Application Tag
        _add_textbox(
            slide,
            left_pos + Inches(0.3),
            Inches(2.1),
            col_width - Inches(0.6),
            Inches(0.3),
            f"{get_text('usecase_label', lang)} 0{idx + 1}",
            font_size=11,
            bold=True,
            color=ACCENT_PINK,
            lang=lang,
        )
        # Application Name
        _add_textbox(
            slide,
            left_pos + Inches(0.3),
            Inches(2.5),
            col_width - Inches(0.6),
            Inches(0.8),
            app.get("application", "Application"),
            font_size=18,
            bold=True,
            color=TEXT_PRIMARY,
            lang=lang,
        )
        # Description
        _add_textbox(
            slide,
            left_pos + Inches(0.3),
            Inches(3.5),
            col_width - Inches(0.6),
            Inches(2.8),
            app.get("description", "Description"),
            font_size=13,
            color=TEXT_SECONDARY,
            line_spacing=18,
            lang=lang,
        )

    # ── Slide 7: Future Outlook & Recommendations ────────────────────
    slide = _create_slide_with_header(
        prs,
        get_text("sec_outlook", lang),
        get_text("strategic_predictions_heading", lang),
        lang=lang,
    )

    # Left Card: Future Outlook / Predictions
    _add_card(slide, Inches(0.8), Inches(1.8), Inches(7.5), Inches(4.8))
    _add_textbox(
        slide,
        Inches(1.1),
        Inches(2.1),
        Inches(6.9),
        Inches(0.4),
        get_text("sec_future_roadmap", lang),
        font_size=12,
        bold=True,
        color=ACCENT_PINK,
        lang=lang,
    )

    tx_box = slide.shapes.add_textbox(
        Inches(1.1), Inches(2.6), Inches(6.9), Inches(3.8)
    )
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    for idx, outlook in enumerate(report.future_outlook[:4]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"▸  {outlook}"
        p.font.name = _get_font_name(lang)
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(12)
        p.line_spacing = Pt(18)

    # Right Card: References / Citations Sidebar
    _add_card(slide, Inches(8.6), Inches(1.8), Inches(3.933), Inches(4.8))
    _add_textbox(
        slide,
        Inches(8.9),
        Inches(2.1),
        Inches(3.3),
        Inches(0.4),
        get_text("sources_references", lang),
        font_size=12,
        bold=True,
        color=ACCENT_PINK,
        lang=lang,
    )

    ref_tx_box = slide.shapes.add_textbox(
        Inches(8.9), Inches(2.6), Inches(3.3), Inches(3.8)
    )
    ref_tf = ref_tx_box.text_frame
    ref_tf.word_wrap = True
    ref_tf.margin_left = Inches(0.05)
    ref_tf.margin_right = Inches(0.05)
    ref_tf.margin_top = Inches(0.05)
    ref_tf.margin_bottom = Inches(0.05)
    for idx, ref in enumerate(report.references[:4]):
        p_title = ref_tf.paragraphs[0] if idx == 0 else ref_tf.add_paragraph()
        p_title.text = f"{idx + 1}. {ref.get('title', 'Source')}"
        p_title.font.name = _get_font_name(lang)
        p_title.font.size = Pt(11)
        p_title.font.bold = True
        p_title.font.color.rgb = ACCENT_VIOLET
        p_title.space_after = Pt(1)

        p_url = ref_tf.add_paragraph()
        url_str = ref.get("url", "")
        domain = re.sub(r"https?://(www\.)?", "", url_str).split("/")[0]
        p_url.text = f"   Source: {domain}"
        p_url.font.name = _get_font_name(lang)
        p_url.font.size = Pt(9)
        p_url.font.color.rgb = TEXT_SECONDARY
        p_url.space_after = Pt(8)

    # ── Save Presentation ────────────────────────────────────────────
    prs.save(str(filepath))
    logger.info("Premium PPTX generated in language '%s': %s", lang, filepath)
    return str(filepath)
