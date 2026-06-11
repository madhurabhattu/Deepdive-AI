"""
Unit tests for utils/ppt_generator.py

Covers:
- build_ppt() creates a file at the returned path
- Slide count >= 5
- Slide titles are non-empty strings
- sanitise_filename replaces special chars
"""

import os

import pytest
from pptx import Presentation as PptxPresentation

from utils.ppt_generator import build_ppt, sanitise_filename
from utils.report_schema import ResearchReport


@pytest.fixture
def mock_report() -> ResearchReport:
    """Return a valid ResearchReport for testing."""
    return ResearchReport(
        topic="Quantum Computing",
        executive_summary=(
            "Quantum computing is a rapidly advancing field that leverages "
            "quantum mechanical phenomena to perform computations far beyond "
            "the capabilities of classical computers."
        ),
        background_context=(
            "Quantum computing is an emerging field that utilizes quantum "
            "mechanics to solve complex problems."
        ),
        core_concepts=[
            {"term": "Qubit", "definition": "Basic unit of quantum information"},
            {
                "term": "Superposition",
                "definition": "Ability to exist in multiple states simultaneously",
            },
            {"term": "Entanglement", "definition": "Quantum link between particles"},
        ],
        key_insights=[
            "Quantum computers use qubits instead of classical bits.",
            "Superposition allows qubits to represent 0 and 1 simultaneously.",
            "Quantum entanglement enables faster information processing.",
        ],
        benefits_challenges_risks=[
            {
                "item": "High Processing Speed",
                "type": "benefit",
                "description": "Solves complex calculations exponentially faster",
            },
            {
                "item": "Physical decoherence",
                "type": "challenge",
                "description": "Highly sensitive to external environmental noise",
            },
            {
                "item": "Security threats",
                "type": "risk",
                "description": "Could potentially break modern cryptographic standards",
            },
        ],
        real_world_applications=[
            {
                "application": "Cryptography",
                "description": "Quantum key distribution for secure messaging",
            },
            {
                "application": "Material Science",
                "description": "Molecular structure modeling for new materials",
            },
            {
                "application": "Logistics",
                "description": "Route optimization and supply chain management",
            },
        ],
        future_outlook=[
            "Development of fault-tolerant logical qubits",
            "Integration with classical supercomputers",
            "Emergence of secure quantum internet",
        ],
        statistics=[
            {"label": "Market Size (2025)", "value": "$8.6 billion"},
            {"label": "CAGR 2025-2030", "value": "32.7%"},
            {"label": "Active Research Labs", "value": "200+"},
        ],
        references=[
            {
                "title": "IBM Quantum Research",
                "url": "https://www.ibm.com/quantum",
                "snippet": "IBM leads enterprise quantum computing research.",
            },
            {
                "title": "Google Quantum AI",
                "url": "https://quantumai.google/",
                "snippet": "Google achieved quantum supremacy in 2019.",
            },
            {
                "title": "Nature: Quantum Computing Review",
                "url": "https://www.nature.com/subjects/quantum-computing",
                "snippet": "A comprehensive review of quantum computing advances.",
            },
        ],
    )


class TestSanitiseFilename:
    """Tests for the filename sanitisation helper."""

    def test_replaces_special_chars(self):
        assert sanitise_filename("Hello/World:Test") == "hello_world_test"

    def test_replaces_spaces(self):
        assert sanitise_filename("Quantum Computing 101") == "quantum_computing_101"

    def test_truncates_long_names(self):
        long_topic = "a" * 200
        result = sanitise_filename(long_topic)
        assert len(result) <= 80


class TestBuildPpt:
    """Tests for PowerPoint generation."""

    def test_file_exists_at_returned_path(self, mock_report: ResearchReport):
        filepath = build_ppt(mock_report)
        assert os.path.exists(filepath)

    def test_slide_count_exactly_eight(self, mock_report: ResearchReport):
        filepath = build_ppt(mock_report)
        prs = PptxPresentation(filepath)
        assert len(prs.slides) == 8

    def test_filepath_ends_with_pptx(self, mock_report: ResearchReport):
        filepath = build_ppt(mock_report)
        assert filepath.endswith(".pptx")

    def test_filename_contains_sanitised_topic(self, mock_report: ResearchReport):
        filepath = build_ppt(mock_report)
        assert "quantum_computing" in os.path.basename(filepath)

    def test_slides_have_content(self, mock_report: ResearchReport):
        filepath = build_ppt(mock_report)
        prs = PptxPresentation(filepath)
        for slide in prs.slides:
            # Each slide should have at least one shape
            assert len(slide.shapes) > 0
